import os
import json
import sys
import shutil
from dotenv import load_dotenv
from pathlib import Path
from groq import Groq

from memory import LongTermMemoryStore
from sandbox import AgentSandbox

root_dir = Path(__file__).resolve().parent
load_dotenv(dotenv_path=root_dir / ".env")

sys.path.append(os.path.dirname(os.path.abspath(__file__))) 

groq_key = os.getenv("GROQ_API_KEY")

if groq_key and groq_key.strip():
    client = Groq(api_key=groq_key)
    MODEL_NAME = "openai/gpt-oss-120b"
else:
    print("Critical System Error: GROQ_API_KEY was not detected in your workspace .env file.")
    sys.exit(1)


class DeepAgent:
    def __init__(self, objective: str, user_id: str):
        self.objective = objective
        self.max_execution_rounds = 4

        self.memory_store = LongTermMemoryStore()
        self.user_namespace = user_id
        self.sandbox = AgentSandbox(scope="assistant", identifier="global")

        self.skills_rubric = {
            "docx": "Exclusively for .docx targets. Extract text contents via 'read_office_file'. To view layout details or tracked changes, unzip the file container using 'unpack_raw_xml'.",
            "pptx": "Triggers on decks/slides. Pull slide outlines using 'read_office_file'. For visual QA, use subagent verification states.",
            "xlsx": "Handles tabular formats. CRITICAL CONSTRAINT: Enforce formula injection! Never hardcode mathematical summaries or averages using scalar numbers. Write active Excel formulas (e.g., =SUM(B1:B10))."
        }

    def spawn_worker_subagent(self, context_key: str, worker_prompt: str) -> str:
        print(f"Synchronous Subagent for context: {context_key}")
        
        subagent_base_prompt = (
            f"You are a specialized engineering subagent workspace runner.\n"
            f"Your strict domain parameter boundaries are: {self.skills_rubric.get(context_key, '')}\n"
            f"Execute your assigned task diligently and return only your completed findings."
        )

        messages = [
            {"role": "system", "content": subagent_base_prompt},
            {"role": "user", "content": worker_prompt}
        ]
        res = client.chat.completions.create(model=MODEL_NAME, messages=messages, temperature=0.01)
        return res.choices[0].message.content

    
    def execute_react_loop(self):
        historical_context = self.memory_store.get_all_context(self.user_namespace)

        base_system_prompt = f"""You are a master deep agent cluster supervisor. You orchestrate file mutations and guide subagents.
Instead of doing manual work yourself, delegate specific components to subagents or run sandbox tools.

To invoke tool chains inside your sandbox workspace, respond using this JSON structure:
{{
    "action": "call_tool",
    "tool_name": "create_xlsx_with_formulas",
    "args": {{"filepath": "ledger.xlsx", "json_data_matrix": [["Total", "=SUM(B2:B3)"]]}}
}}

To delegate complex research or validation tasks to an isolated subagent:
{{
    "action": "delegate_subagent",
    "skill_context_key": "xlsx",
    "prompt": "Audit the formula layout structure to verify no values are hardcoded."
}}

To execute environment terminal setup operations via your tool setup:
{{
    "action": "call_tool",
    "tool_name": "execute_shell_command",
    "args": {{"command": "ls -la"}}
}}

[LONG-TERM MEMORY ENVIRONMENT CONTEXT]:
{historical_context}

[ACTIVE SKILLS BLUEPRINT REFERENCE]:
{json.dumps(self.skills_rubric)}
"""

        messages = [
            {"role": "system", "content": base_system_prompt},
            {"role": "user", "content": self.objective}
        ]

        for loop_idx in range(1, self.max_execution_rounds + 1):
            print(f"\n Loop Round {loop_idx} out of {self.max_execution_rounds}")
            res = client.chat.completions.create(model=MODEL_NAME, messages=messages, temperature=0.01)
            
            thought = res.choices[0].message.content
            print(f"Agent Thought Process: \n {thought}")
            messages.append({"role": "assistant", "content": thought})

            executed_any = False
            bracket_count = 0
            start_idx = -1
            
            for idx, char in enumerate(thought):
                if char == '{':
                    if bracket_count == 0:
                        start_idx = idx
                    bracket_count += 1
                elif char == '}':
                    bracket_count -= 1
                    if bracket_count == 0 and start_idx != -1:
                        json_string = thought[start_idx:idx+1]
                        try:
                            payload = json.loads(json_string)
                            if "action" in payload:
                                executed_any = True
                                if payload["action"] == "call_tool":
                                    tool_return = self.sandbox.run_sandbox_tool(payload["tool_name"], **payload["args"])
                                    print(f"Tool [{payload['tool_name']}] Output: {tool_return[:150]}")
                                    messages.append({"role": "user", "content": f"Tool Execution Output: {tool_return}"})

                                elif payload["action"] == "delegate_subagent":
                                    sub_return = self.spawn_worker_subagent(payload["skill_context_key"], payload["prompt"])
                                    print(f"Subagent Worker Output:\n{sub_return}")
                                    messages.append({"role": "user", "content": f"Subagent Output: {sub_return}"})
                        except Exception as json_err:
                            print(f"Individual block parsing skipped: {str(json_err)}")
                        start_idx = -1

            if executed_any:
                continue  
            break

        self.memory_store.save_memory(
            self.user_namespace,
            "last_executed_summary",
            {"objective_status": "Success", "handled_path": "ledger.xlsx"}
        )

        self.sandbox.close_and_cleanup()

if __name__ == "__main__":
    test_instruction = (
        "Read the text layout from the presentation 'sales.pptx', use those findings to create a clean "
        "financial spreadsheet workbook named 'ledger.xlsx' enforcing strict Excel formulas for summary cells, "
        "and delegate a subagent to run an audit on the formula architecture variables."
    )
    
    orchestrator = DeepAgent(objective=test_instruction, user_id="user_101")

    local_source_file = Path("sales.pptx")
    if not local_source_file.exists():
        from pptx import Presentation
        prs = Presentation()
        
        slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide_1.shapes.title.text = "Q3 Corporate Revenue Summary"
        
        slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide_2.shapes.title.text = "Financial Targets"
        body_shape = slide_2.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Core Performance Metrics:"
        tf.add_paragraph().text = "- Expected License Sales: 50000"
        tf.add_paragraph().text = "- Projected Cloud Revenue: 120000"
        tf.add_paragraph().text = "- Consulting Services: 30000"
        
        prs.save(str(local_source_file))
        print(f"Created missing data-seeded presentation: {local_source_file}")
    else:
        if local_source_file.stat().st_size < 30000:  
            local_source_file.unlink()
            print("Cleared old empty presentation file cache. Please re-run the script to execute with fresh data seeding.")
            sys.exit(0)

    orchestrator.execute_react_loop()