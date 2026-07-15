import os
import json
from pathlib import Path
from markitdown import MarkItDown
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from docx import Document
from pptx import Presentation
from langchain_core.tools import tool

md_converter = MarkItDown()

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _resolve(filepath: str) -> str:
    """Always write into OUTPUT_DIR, regardless of what path style the model
    used — a real Windows absolute path, a virtual FilesystemBackend '/' path
    like '/outputs/ledger.xlsx', or a bare filename. Only the filename itself
    matters; any directory portion the model supplies is discarded, since
    Path.is_absolute() unreliably treats a leading '/' as absolute on Windows
    even when it's really a virtual-backend path, not a real filesystem path."""
    name = Path(filepath.replace("\\", "/")).name
    if not name:
        name = "output_file"
    resolved = OUTPUT_DIR / name
    resolved.parent.mkdir(parents=True, exist_ok=True)
    return str(resolved)

@tool
def read_office_file(filepath: str) -> str:
    """Extract plain text/markdown content from an office file (.pptx, .docx, or .xlsx)."""
    filepath = _resolve(filepath)
    if not os.path.exists(filepath):
        return f"Error: Target file '{filepath}' not found on the local filesystem."
    try:
        result = md_converter.convert(filepath)
        return result.text_content
    except Exception as e:
        return f"MarkItDown extraction failed: {str(e)}"

@tool
def create_xlsx_with_formulas(filepath: str, json_data_matrix: list) -> dict:
    """Create a new Excel spreadsheet. Pass json_data_matrix as a list of rows; any
    calculated cell must be a formula string starting with '=' (e.g. '=SUM(B2:B4)'),
    never a hardcoded number."""
    filepath = _resolve(filepath)
    try:
        rows = json.loads(json_data_matrix) if isinstance(json_data_matrix, str) else json_data_matrix

        wb = Workbook()
        ws = wb.active
        for row in rows:
            processed_row = [
                cell["formula"] if isinstance(cell, dict) and "formula" in cell else cell
                for cell in row
            ]
            ws.append(processed_row)

        for cell in ws[1]:
            cell.font = Font(bold=True)

        for col_idx, col_cells in enumerate(ws.columns, start=1):
            max_len = max((len(str(c.value)) for c in col_cells if c.value is not None), default=10)
            ws.column_dimensions[get_column_letter(col_idx)].width = max_len + 2

        wb.save(filepath)
        return {"success": True, "message": f"Successfully created Excel spreadsheet at {filepath}.", "path": filepath}
    except Exception as e:
        return {"success": False, "message": f"Spreadsheet build failed: {str(e)}", "path": filepath}

@tool
def create_word_document(filepath: str, content: str, title: str = "Executive Report") -> dict:
    """Create a new Word document with the given title and content. Use only when the
    target file does not already exist. Lines starting with '#' become headings; lines
    starting with '- ' become bullet list items. Content must be primarily narrative
    prose — bullet lines are capped at 35% of total content lines."""
    filepath = _resolve(filepath)
    pre_existing = os.path.exists(filepath)
    try:
        lines = [l for l in str(content).split("\n") if l.strip()]
        bullet_lines = [l for l in lines if l.strip().startswith("- ")]
        bullet_ratio = len(bullet_lines) / len(lines) if lines else 0

        if bullet_ratio > 0.35:
            return {
                "success": False,
                "message": (
                    f"Rejected: {len(bullet_lines)} of {len(lines)} lines ({bullet_ratio:.0%}) "
                    "are bullet points, exceeding the `35%` limit. Rewrite this content as "
                    "connected narrative paragraphs. Call this tool again with revised content."
                ),
                "path": filepath,
            }

        doc = Document()
        doc.add_heading(title, level=0)
        for line in lines:
            if line.strip().startswith("#"):
                doc.add_heading(line.replace("#", "").strip(), level=1)
            elif line.strip().startswith("- "):
                doc.add_paragraph(line.strip()[2:], style="List Bullet")
            else:
                doc.add_paragraph(line)

        doc.save(filepath)
        return {"success": True, "message": f"{'Overwrote existing' if pre_existing else 'Created new'} Word Document at {filepath}.", "path": filepath, "pre_existing": pre_existing,}
    except Exception as e:
        return {"success": False, "message": f"Word document creation failed: {str(e)}", "path": filepath}

@tool
def append_text_to_document(filepath: str, content: str) -> dict:
    """Append one paragraph of content to an existing Word document, creating it first
    if it doesn't exist yet."""
    filepath = _resolve(filepath)
    try:
        if not os.path.exists(filepath):
            doc = Document()
            doc.add_heading("Automated Updates Document log", level=1)
        else:
            doc = Document(filepath)

        doc.add_paragraph(content)
        doc.save(filepath)
        return {"success": True, "message": f"Appended content updates to {filepath}.", "path": filepath}
    except Exception as e:
        return {"success": False, "message": f"Failed to append content onto text file target: {str(e)}", "path": filepath}

@tool
def modify_presentation_metadata(filepath: str, action_type: str, raw_text_content: str, slide_title: str = "") -> dict:
    """Add a new slide to a PowerPoint file (creating the file if missing). Only
    action_type='add_slide' is currently supported.

    slide_title: the specific, descriptive title for THIS slide (e.g. 'Average Lifespan',
    not a generic label). Every slide must have a distinct title reflecting its own content.

    IMPORTANT: there is no way to edit or delete a slide once added. Before your first call,
    plan your complete list of slides and their titles. Call this tool exactly once per
    planned slide — never speculatively, and never to redo a slide you already added.
    Check the returned total_slides count after each call to track your progress."""
    filepath = _resolve(filepath)
    try:
        pres = Presentation(filepath) if os.path.exists(filepath) else Presentation()
        if action_type == "add_slide":
            slide = pres.slides.add_slide(pres.slide_layouts[1])
            lines = [l for l in raw_text_content.split("\n") if l.strip()]
            title = slide_title.strip() or (lines[0] if lines else "Untitled Slide")
            slide.shapes.title.text = title
            body_lines = lines[1:] if lines and lines[0] == title else lines
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = body_lines[0] if body_lines else ""
            for line in body_lines[1:]:
                tf.add_paragraph().text = line
        pres.save(filepath)
        total = len(pres.slides)
        return {
            "success": True,
            "message": f"Added slide '{title}'. The file now has {total} slide(s) total.",
            "path": filepath,
            "total_slides": total,
        }
    except Exception as e:
        return {"success": False, "message": f"Presentation pipeline modifications failed: {str(e)}", "path": filepath}
    

ALL_TOOLS = [
    read_office_file,
    create_xlsx_with_formulas,
    create_word_document,
    append_text_to_document,
    modify_presentation_metadata,
]
TOOLS_BY_NAME = {t.name: t for t in ALL_TOOLS}