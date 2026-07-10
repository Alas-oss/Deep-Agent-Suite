import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
sys.stderr.reconfigure(encoding="utf-8", errors="replace")
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent
SRC_DIR = REPO_ROOT / "src"
OUTPUT_DIR = REPO_ROOT / "outputs"
sys.path.insert(0, str(SRC_DIR))

from agent import DeepAgent, langfuse
from scenarios import SCENARIOS
from tui import print_banner, console

def main():
    if len(sys.argv) > 1 and sys.argv[1].isdigit():
        _run_single(SCENARIOS[int(sys.argv[1])], user_id=_default_user_id())
        return
    if len(sys.argv) > 1:
        _run_single({"name": "custom", "prompt": " ".join(sys.argv[1:]), "expected_outputs": []}, user_id=_default_user_id())
        return
    
    _ensure_seed_files()
    print_banner("Deep Agent Suite - interactive session")
    print("Type a task for the agent, or 'quit' to exit.")
    print("Supported areas: reading/summarizing .pptx/.docx/.xlsx files, building spreadsheets "
          "with live formulas, writing Word reports, and adding PowerPoint slides.\n")
    
    typed_id = console.input(f"Session name (for '{_default_user_id()}'): ").strip()
    session_user_id = typed_id or _default_user_id()

    while True:
        user_prompt = console.input("> ").strip()
        if not user_prompt or user_prompt.lower() in ("quit", "exit"):
            print("Session ended")
            break

        console.print(f"\nRunning: {user_prompt}\n")

        orchestrator = DeepAgent(objective=user_prompt, user_id=session_user_id) 
        try: 
            orchestrator.execute_react_loop()
        finally:
            langfuse.flush()
        print()

def _default_user_id() -> str:
    import getpass
    try:
        return getpass.getuser() or "anonymous"
    except Exception:
        return "anonymous"

def _run_single(scenario, user_id="anonymous"):
    _ensure_seed_files()
    print_banner(f"Running scenario: {scenario['name']}")
    orchestrator = DeepAgent(objective=scenario["prompt"], user_id=user_id)
    try:
        orchestrator.execute_react_loop()
    finally:
        langfuse.flush()

def _ensure_seed_files():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    _ensure_seed_pptx()
    _ensure_seed_docx()
    _ensure_seed_xlsx()


def _ensure_seed_pptx():
    seed_path = OUTPUT_DIR / "sales.pptx"
    if seed_path.exists():
        return
    #try to get it to create an actual presentation when needed
    #Can be done with a template or smth of the sort
    from pptx import Presentation
    pres = Presentation()
    slide_1 = pres.slides.add_slide(pres.slide_layouts[0])
    slide_1.shapes.title.text = "Q3 Corporate Revenue Summary"
    slide_2 = pres.slides.add_slide(pres.slide_layouts[1])
    slide_2.shapes.title.text = "Financial Targets"
    body_shape = slide_2.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Core Performance Metrics:"
    tf.add_paragraph().text = "- Expected License Sales: 50000"
    tf.add_paragraph().text = "- Projected Cloud Revenue: 120000"
    tf.add_paragraph().text = "- Consulting Services: 30000"
    pres.save(str(seed_path))
    print(f"Created missing data-seeded presentation: {seed_path}")

def _ensure_seed_docx():
    seed_path = OUTPUT_DIR / "executive_report.docx"
    if seed_path.exists() and _docx_is_valid(seed_path):
        return
    
    from docx import Document
    doc = Document()
    doc.add_heading("Executive Report", level=0)
    doc.add_paragraph("Placeholder report - populated by earlier scenario runs.")
    doc.save(str(seed_path))
    print(f"Created missing placeholder report: {seed_path}")

def _docx_is_valid(path) -> bool:
    from docx import Document
    try:
        Document(str(path))
        return True
    except Exception:
        return False

def _ensure_seed_xlsx():
    seed_path = OUTPUT_DIR / "q3_metrics.xlsx"
    if seed_path.exists():
        return

    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["Metric", "Value"])
    ws.append(["License Sales", 50000])
    ws.append(["Cloud Revenue", 120000])
    ws.append(["Consulting Services", 30000])
    ws.append(["Total", "=SUM(B2:B4)"])
    wb.save(str(seed_path))
    print(f"Created missing data-seeded spreadsheet: {seed_path}")

if __name__=="__main__":
    main()